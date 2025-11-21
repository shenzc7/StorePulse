#!/usr/bin/env python3
"""Generate a polished StorePulse PowerPoint deck."""

import os
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --- Brand palette ---------------------------------------------------------
PRIMARY_BLUE = RGBColor(28, 64, 147)
DEEP_NAVY = RGBColor(16, 24, 48)
MID_BLUE = RGBColor(64, 115, 255)
SKY_BLUE = RGBColor(225, 235, 255)
AQUA = RGBColor(0, 178, 169)
ENERGY = RGBColor(255, 140, 66)
PURPLE = RGBColor(133, 97, 241)
LIME = RGBColor(120, 190, 32)
STEEL = RGBColor(90, 104, 128)
COOL_GREY = RGBColor(239, 243, 250)
WHITE = RGBColor(255, 255, 255)

FONT_PRIMARY = "Calibri"


# --- Helper utilities ------------------------------------------------------
def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COOL_GREY


def add_kicker(slide, text):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(8), Inches(0.4))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(16)
    p.font.color.rgb = STEEL
    p.font.name = FONT_PRIMARY
    p.font.bold = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    return box


def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.95), Inches(8.2), Inches(1.2))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.name = FONT_PRIMARY
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE
    frame.vertical_anchor = MSO_ANCHOR.TOP
    return title_box


def add_subtitle(slide, text, left=0.9, top=1.6, width=8.2, height=1.2):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.name = FONT_PRIMARY
    p.font.color.rgb = STEEL
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.word_wrap = True
    return box


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.7), Inches(0.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.solid()
    bar.line.width = Pt(0)
    return bar


def add_card(slide, title, body, left, top, width, height, fill_color, text_color, title_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.width = Pt(0)
    card.shadow.inherit = False

    frame = card.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(16)
    frame.margin_right = Pt(16)
    frame.margin_top = Pt(16)

    header = frame.paragraphs[0]
    header.text = title
    header.font.size = Pt(20)
    header.font.bold = True
    header.font.name = FONT_PRIMARY
    header.font.color.rgb = title_color or text_color

    body_paragraph = frame.add_paragraph()
    body_paragraph.text = body
    body_paragraph.font.size = Pt(14)
    body_paragraph.font.name = FONT_PRIMARY
    body_paragraph.font.color.rgb = text_color
    body_paragraph.space_before = Pt(6)
    return card


def add_checklist(slide, items, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(4)
    frame.word_wrap = True

    for idx, text in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = f"• {text}"
        paragraph.font.size = Pt(18 if idx == 0 else 16)
        paragraph.font.name = FONT_PRIMARY
        paragraph.font.color.rgb = DEEP_NAVY
        if idx == 0:
            paragraph.font.bold = True
    return box


def add_timeline(slide, steps, left, top, spacing):
    for idx, (title, subtitle) in enumerate(steps):
        x = left + idx * spacing
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(top), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = MID_BLUE if idx == 0 else PRIMARY_BLUE
        circle.line.width = Pt(0)

        label_box = slide.shapes.add_textbox(Inches(x - 0.25), Inches(top + 0.6), Inches(1), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = str(idx + 1)
        label_frame.paragraphs[0].font.size = Pt(14)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(x - 0.45), Inches(top + 1.0), Inches(1.4), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(16)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = DEEP_NAVY
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True

        subtitle_box = slide.shapes.add_textbox(Inches(x - 0.65), Inches(top + 1.55), Inches(1.8), Inches(0.9))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(13)
        subtitle_frame.paragraphs[0].font.color.rgb = STEEL
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.word_wrap = True

        if idx < len(steps) - 1:
            connector = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.5),
                Inches(top + 0.22),
                Inches(spacing - 0.1),
                Inches(0.05),
            )
            connector.fill.solid()
            connector.fill.fore_color.rgb = SKY_BLUE
            connector.line.width = Pt(0)


# --- Slide constructors ----------------------------------------------------
def add_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Hero backdrop
    hero = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.6),
        Inches(8.8),
        Inches(4.9),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = DEEP_NAVY
    hero.line.width = Pt(0)

    overlay = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.8),
        Inches(0.6),
        Inches(3.6),
        Inches(4.9),
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = MID_BLUE
    overlay.line.width = Pt(0)

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(4.8), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True

    title = title_frame.paragraphs[0]
    title.text = "StorePulse"
    title.font.size = Pt(56)
    title.font.bold = True
    title.font.name = FONT_PRIMARY
    title.font.color.rgb = WHITE

    subtitle = title_frame.add_paragraph()
    subtitle.text = "AI-powered retail demand foresight"
    subtitle.font.size = Pt(24)
    subtitle.font.name = FONT_PRIMARY
    subtitle.font.color.rgb = SKY_BLUE

    tagline_box = slide.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(4.6), Inches(1.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = (
        "MTech Capstone • Computer Science + Business Analytics\n"
        "Know tomorrow's visits. Act today."
    )
    tagline_frame.paragraphs[0].font.size = Pt(18)
    tagline_frame.paragraphs[0].font.color.rgb = WHITE
    tagline_frame.paragraphs[0].font.name = FONT_PRIMARY
    tagline_frame.word_wrap = True

    # Highlight cards
    highlight_data = [
        ("14-day demand horizon", "Confidence bands + action plans"),
        ("Offline-first desktop", "macOS • Windows • ₹0 infra"),
        ("Actionable intelligence", "Staffing + inventory recommendations"),
    ]

    for idx, (title_text, body_text) in enumerate(highlight_data):
        card_left = 1.1 + idx * 3.0
        add_card(
            slide,
            title_text,
            body_text,
            left=card_left,
            top=3.6,
            width=2.7,
            height=1.0,
            fill_color=WHITE,
            text_color=DEEP_NAVY,
            title_color=PRIMARY_BLUE,
        )

    return slide


def add_vision_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Context")
    add_title(slide, "Why retailers need StorePulse now")
    add_accent_bar(slide)

    left_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(4.6), Inches(3.6))
    frame = left_box.text_frame
    frame.clear()
    frame.word_wrap = True

    bullet_points = [
        "Footfall volatility makes weekly heuristics unreliable",
        "Store managers lack time to interpret raw analytics",
        "Business teams demand traceable, defensible forecasts",
    ]

    for idx, text in enumerate(bullet_points):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = text
        paragraph.font.size = Pt(20 if idx == 0 else 18)
        paragraph.font.name = FONT_PRIMARY
        paragraph.font.color.rgb = DEEP_NAVY
        paragraph.space_after = Pt(10)
        paragraph.level = 0

    insight_cards = [
        ("18-25%", "Average staffing mismatch caused by manual scheduling"),
        ("3x", "Higher promo ROI when forecasts power inventory"),
        ("<90 sec", "From fresh data to an actionable StorePulse plan"),
    ]

    for idx, (title_text, body_text) in enumerate(insight_cards):
        add_card(
            slide,
            title_text,
            body_text,
            left=5.3,
            top=2.0 + idx * 1.4,
            width=3.0,
            height=1.2,
            fill_color=PRIMARY_BLUE if idx != 2 else AQUA,
            text_color=WHITE,
        )

    return slide


def add_pain_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Problem")
    add_title(slide, "Where current retail planning breaks")
    add_accent_bar(slide)

    cards = [
        (
            "Blind spots",
            "Spreadsheets ignore holidays, weather swings, and local events leading to whiplash staffing decisions.",
            PURPLE,
        ),
        (
            "Late reactions",
            "Teams notice demand shifts only after lost sales or overtime burns budget.",
            ENERGY,
        ),
        (
            "Data friction",
            "Store associates cannot wrangle CSVs or large BI tools before doors open.",
            PRIMARY_BLUE,
        ),
    ]

    for idx, (title_text, body_text, color) in enumerate(cards):
        add_card(
            slide,
            title_text,
            body_text,
            left=0.9 + idx * 3.0,
            top=2.0,
            width=2.7,
            height=3.2,
            fill_color=color,
            text_color=WHITE,
        )

    return slide


def add_solution_snapshot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Solution Overview")
    add_title(slide, "StorePulse in one slide")
    add_accent_bar(slide)

    cards = [
        (
            "Forecasting Engine",
            "Negative-Binomial ARX with Bayesian calibration delivers trustworthy visit predictions and P10/P50/P90 bands.",
        ),
        (
            "Action Planner",
            "Turns forecasts into shift rosters and SKU-level stock moves, ready to export as PDF or WhatsApp-friendly summaries.",
        ),
        (
            "Insight Hub",
            "Scenario lab to test promotions, weather shocks, or staffing constraints before committing to the floor plan.",
        ),
    ]

    for idx, (title_text, body_text) in enumerate(cards):
        add_card(
            slide,
            title_text,
            body_text,
            left=0.9,
            top=2.0 + idx * 1.5,
            width=8.0,
            height=1.3,
            fill_color=WHITE,
            text_color=DEEP_NAVY,
            title_color=PRIMARY_BLUE,
        )

    return slide


def add_modes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Product Design")
    add_title(slide, "Progressive adoption: Lite to Pro")
    add_accent_bar(slide)

    lite_card = add_card(
        slide,
        "Lite Mode",
        "• Input: date + visits\n• Visual wizard in <2 minutes\n• Auto-learns weekday and holiday patterns\n• 8%+ accuracy lift vs 7-day average",
        left=0.9,
        top=2.0,
        width=4.0,
        height=3.2,
        fill_color=WHITE,
        text_color=DEEP_NAVY,
        title_color=AQUA,
    )

    pro_card = add_card(
        slide,
        "Pro Mode",
        "• Adds promotions, pricing, weather, events\n• LightGBM residual booster for peak sensitivity\n• 20%+ weekend error reduction\n• Detailed staffing + SKU adjustments",
        left=5.1,
        top=2.0,
        width=3.8,
        height=3.2,
        fill_color=WHITE,
        text_color=DEEP_NAVY,
        title_color=PRIMARY_BLUE,
    )

    for card in (lite_card, pro_card):
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = SKY_BLUE
        card.line.width = Pt(2)

    return slide


def add_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Architecture")
    add_title(slide, "Experience • Intelligence • Platform")
    add_accent_bar(slide)

    layers = [
        (
            "Experience Layer",
            "Tauri + React + Tailwind\nRadix UI patterns\nOffline-first desktop packaging",
            WHITE,
        ),
        (
            "Intelligence Layer",
            "NB-ARX core model\nBayesian uncertainty module\nLightGBM booster + conformal calibration",
            WHITE,
        ),
        (
            "Platform Layer",
            "FastAPI services\nSQLite local data lake\nTask orchestration, PDF/CSV exporters",
            WHITE,
        ),
    ]

    for idx, (title_text, body_text, fill_color) in enumerate(layers):
        add_card(
            slide,
            title_text,
            body_text,
            left=0.9,
            top=2.0 + idx * 1.5,
            width=8.0,
            height=1.3,
            fill_color=fill_color,
            text_color=DEEP_NAVY,
            title_color=PRIMARY_BLUE,
        )

    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.0), Inches(0.2), Inches(3.9))
    connector.fill.solid()
    connector.fill.fore_color.rgb = MID_BLUE
    connector.line.width = Pt(0)

    return slide


def add_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "ML Flow")
    add_title(slide, "Trusted pipeline from data to decisions")
    add_accent_bar(slide)

    steps = [
        ("Capture", "CSV upload, manual entry, POS sync"),
        ("Feature", "Holidays, weather, promos, paydays"),
        ("Model", "NB-ARX + Bayesian NB + LightGBM"),
        ("Calibrate", "Inductive conformal coverage"),
        ("Act", "Shift rosters, SKU moves, PDF export"),
    ]

    add_timeline(slide, steps, left=1.1, top=2.1, spacing=1.6)
    return slide


def add_engine_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Forecasting Engine")
    add_title(slide, "Bayesian accuracy with interpretable levers")
    add_accent_bar(slide)

    checklist_items = [
        "Negative-Binomial ARX handles count data + exogenous drivers",
        "Bayesian posterior draws deliver honest confidence bands",
        "LightGBM residual learner boosts promo + weather spikes",
        "Explainable contribution charts clarify drivers for managers",
    ]
    add_checklist(slide, checklist_items, left=0.9, top=2.05, width=4.0, height=3.2)

    chart_data = CategoryChartData()
    chart_data.categories = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]
    chart_data.add_series("Actual", (210, 240, 205, 260, 320, 402, 376))
    chart_data.add_series("Forecast", (208, 236, 212, 255, 315, 395, 368))
    chart_data.add_series("P90", (230, 258, 232, 276, 339, 422, 399))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(5.2),
        Inches(2.1),
        Inches(3.8),
        Inches(3.0),
        chart_data
    ).chart

    chart.chart_title.text_frame.text = "Illustrative weekly forecast"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.has_data_labels = False
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(11)
    value_axis = chart.value_axis
    value_axis.visible = True
    value_axis.tick_labels.font.size = Pt(11)

    return slide


def add_experience_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "User Experience")
    add_title(slide, "Crafted for managers, loved by analysts")
    add_accent_bar(slide)

    cards = [
        (
            "Guided onboarding",
            "Two-screen wizard for Lite mode. Drag-and-drop CSV import with validation for Pro mode.",
            WHITE,
            PRIMARY_BLUE,
        ),
        (
            "What-if lab",
            "Drag sliders for weather, promo intensity, or staffing caps to preview impact before committing.",
            WHITE,
            PRIMARY_BLUE,
        ),
        (
            "Instant sharing",
            "One-click PDF and WhatsApp summaries with staffing rosters and stock deltas.",
            WHITE,
            PRIMARY_BLUE,
        ),
    ]

    for idx, (title_text, body_text, fill_color, title_color) in enumerate(cards):
        add_card(
            slide,
            title_text,
            body_text,
            left=0.9 + idx * 3.0,
            top=2.0,
            width=2.7,
            height=3.2,
            fill_color=fill_color,
            text_color=DEEP_NAVY,
            title_color=title_color,
        )

    return slide


def add_value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Business Impact")
    add_title(slide, "Measurable value across the floor")
    add_accent_bar(slide)

    metrics = [
        ("15-25%", "Reduction in overtime and idle staffing hours"),
        ("12%", "Lower stockouts through proactive replenishment"),
        ("+6 pts", "CSAT lift when visit surges are anticipated"),
    ]

    for idx, (metric, copy) in enumerate(metrics):
        add_card(
            slide,
            metric,
            copy,
            left=0.9 + idx * 3.0,
            top=2.0,
            width=2.7,
            height=1.5,
            fill_color=PRIMARY_BLUE if idx == 0 else (PURPLE if idx == 1 else AQUA),
            text_color=WHITE,
        )

    add_subtitle(
        slide,
        "Retail leadership views StorePulse as an execution engine: something their teams can run daily without data science support.",
        left=0.9,
        top=3.8,
        width=8.0,
        height=1.0,
    )

    return slide


def add_quality_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Quality & Trust")
    add_title(slide, "Governance baked into every build")
    add_accent_bar(slide)

    gates = [
        "Lite mode beats 7-day moving average by ≥8% sMAPE",
        "Pro mode improves weekend accuracy by ≥20% vs Lite",
        "Conformal intervals deliver 80-95% coverage",
        "Cold start to first plan ≤90 seconds on sample data",
        "Automated regression tests guard forecasting + UI flows",
    ]

    add_checklist(slide, gates, left=0.9, top=2.0, width=4.2, height=3.2)

    doc_card = add_card(
        slide,
        "Executive assurance",
        "Quality reports auto-generate with every release, summarizing gate performance for stakeholders.",
        left=5.1,
        top=2.0,
        width=3.8,
        height=1.6,
        fill_color=WHITE,
        text_color=DEEP_NAVY,
        title_color=PRIMARY_BLUE,
    )
    doc_card.line.fill.solid()
    doc_card.line.fill.fore_color.rgb = SKY_BLUE
    doc_card.line.width = Pt(2)

    return slide


def add_demo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Live Demo")
    add_title(slide, "From raw data to action in minutes")
    add_accent_bar(slide)

    steps = [
        ("Ingest", "Import CSV or hand-log yesterday's visits"),
        ("Model", "Select Lite/Pro, run training"),
        ("Forecast", "Review 14-day P10/P50/P90 curves"),
        ("Plan", "Generate staffing + SKU actions"),
        ("Export", "Share PDF + WhatsApp summary"),
    ]

    add_timeline(slide, steps, left=1.1, top=2.1, spacing=1.6)
    return slide


def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_kicker(slide, "Roadmap")
    add_title(slide, "Scaling impact beyond the pilot")
    add_accent_bar(slide)

    short_term = [
        "Portfolio dashboard across stores",
        "Scenario packs for festive seasons",
        "POS connectors (Vend, Shopify)",
    ]
    long_term = [
        "Cloud sync optionality for HQ oversight",
        "Mobile companion for area managers",
        "Marketplace of specialized forecasting plugins",
    ]

    add_card(
        slide,
        "Next 3-6 months",
        "\n".join(f"• {item}" for item in short_term),
        left=0.9,
        top=2.0,
        width=3.8,
        height=3.0,
        fill_color=WHITE,
        text_color=DEEP_NAVY,
        title_color=AQUA,
    )

    add_card(
        slide,
        "6-12 months vision",
        "\n".join(f"• {item}" for item in long_term),
        left=5.1,
        top=2.0,
        width=3.8,
        height=3.0,
        fill_color=WHITE,
        text_color=DEEP_NAVY,
        title_color=PRIMARY_BLUE,
    )

    return slide


def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    closing_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(0.9),
        Inches(8.4),
        Inches(4.2),
    )
    closing_card.fill.solid()
    closing_card.fill.fore_color.rgb = PRIMARY_BLUE
    closing_card.line.width = Pt(0)

    frame = closing_card.text_frame
    frame.clear()
    frame.margin_left = Pt(28)
    frame.margin_top = Pt(28)
    frame.word_wrap = True

    title = frame.paragraphs[0]
    title.text = "Thank you"
    title.font.size = Pt(40)
    title.font.bold = True
    title.font.color.rgb = WHITE
    title.font.name = FONT_PRIMARY

    subtitle = frame.add_paragraph()
    subtitle.text = "StorePulse — reliable retail planning, zero guesswork"
    subtitle.font.size = Pt(20)
    subtitle.font.color.rgb = SKY_BLUE
    subtitle.font.name = FONT_PRIMARY
    subtitle.space_before = Pt(10)

    contact = frame.add_paragraph()
    contact.text = (
        "MTech Project Showcase\n"
        "Computer Science & Business Analytics\n\n"
        "Tech Stack: Python • FastAPI • Tauri • React • PyMC • LightGBM\n"
        "Ready for demo and technical deep-dive"
    )
    contact.font.size = Pt(16)
    contact.font.color.rgb = WHITE
    contact.font.name = FONT_PRIMARY
    contact.space_before = Pt(18)

    return slide


# --- Assembly --------------------------------------------------------------
def create_presentation():
    prs = Presentation()

    add_cover_slide(prs)
    add_vision_slide(prs)
    add_pain_slide(prs)
    add_solution_snapshot(prs)
    add_modes_slide(prs)
    add_architecture_slide(prs)
    add_pipeline_slide(prs)
    add_engine_slide(prs)
    add_experience_slide(prs)
    add_value_slide(prs)
    add_quality_slide(prs)
    add_demo_slide(prs)
    add_roadmap_slide(prs)
    add_closing_slide(prs)

    return prs


def main():
    prs = create_presentation()
    output_path = os.path.join(
        "/Users/shenzc/Desktop/projects/StorePulse",
        "StorePulse_Presentation_MTech.pptx",
    )
    prs.save(output_path)

    print("✅ Presentation created")
    print(f"📁 Saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"📏 File size: {size_mb:.2f} MB")
    print(f"🕒 Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")


if __name__ == "__main__":
    main()
